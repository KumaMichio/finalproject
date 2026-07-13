# -*- coding: utf-8 -*-
"""
build_slides.py — Sinh slide báo cáo ĐATN theo template HUST (.pptx).
Nền = trang template HUST đã render (component_for_final/hust_tpl/), overlay text + ảnh.
Nội dung bám DoAn.pdf. Chạy:
  custom_tracking_system/venv_tracking/Scripts/python custom_tracking_system/scripts/build_slides.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[2]
TPL = ROOT / "component_for_final" / "hust_tpl"
FIG = ROOT / "component_for_final" / "figures_results"
SHOT = ROOT / "component_for_final" / "screenshots"
UML = ROOT / "component_for_final" / "uml" / "png"
AS = ROOT / "docs" / "assets"
OUT = ROOT / "Bao_cao_do_an_tot_nghiep" / "slide_bao_cao.pptx"

NAVY = RGBColor(0x0B, 0x28, 0x50)
RED = RGBColor(0xC1, 0x0E, 0x2B)
DARK = RGBColor(0x1A, 0x22, 0x30)
GRAY = RGBColor(0x55, 0x62, 0x72)
GREEN = RGBColor(0x12, 0x80, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"

prs = Presentation()
prs.slide_width = I(13.333)
prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5
_pageno = [0]


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(str(bg), 0, 0, I(SW), I(SH))
    return s


def rect(s, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(l), I(t), I(w), I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = b.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color=DARK, bold=False, italic=False, bullet=False,
         align=PP_ALIGN.LEFT, space=4, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space); p.space_before = Pt(0)
    r = p.add_run(); r.text = ("•  " + text) if bullet else text
    f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return p


def title(s, text, sub=None):
    tf = tb(s, 0.55, 1.42, 12.2, 1.0)
    para(tf, text, 27, NAVY, bold=True, first=True, space=0)
    if sub:
        tf2 = tb(s, 0.57, 2.16, 12.0, 0.5)
        para(tf2, sub, 14, GRAY, italic=True, first=True, space=0)


def pageno(s):
    _pageno[0] += 1
    rect(s, 12.35, 7.02, 0.85, 0.42, WHITE)
    tf = tb(s, 12.35, 7.03, 0.8, 0.38, MSO_ANCHOR.MIDDLE)
    para(tf, str(_pageno[0]), 13, NAVY, bold=True, align=PP_ALIGN.RIGHT, first=True, space=0)


def content(t, sub=None):
    s = slide(TPL / "pg-05.png")
    title(s, t, sub)
    pageno(s)
    return s


def pic_fit(s, path, l, t, w, h, caption=None):
    im = Image.open(str(path)); iw, ih = im.size
    ar = iw / ih; box = w / h
    if ar > box:
        nw = w; nh = w / ar; nl = l; nt = t + (h - nh) / 2
    else:
        nh = h; nw = h * ar; nt = t; nl = l + (w - nw) / 2
    s.shapes.add_picture(str(path), I(nl), I(nt), I(nw), I(nh))
    if caption:
        tf = tb(s, l, t + h + 0.03, w, 0.35)
        para(tf, caption, 10.5, GRAY, italic=True, align=PP_ALIGN.CENTER, first=True, space=0)


def bullets(s, items, l, t, w, h, size=16, space=7):
    tf = tb(s, l, t, w, h)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        p = para(tf, txt, size - (2 if lvl else 0),
                 DARK if lvl == 0 else GRAY, bullet=(lvl == 0), first=(i == 0), space=space)
        if lvl:
            p.text = "   – " + txt
            for r in p.runs:
                r.font.name = FONT; r.font.size = Pt(size - 2); r.font.color.rgb = GRAY


def kpi(s, items, l, t, w, gap=0.28):
    n = len(items); cw = (w - gap * (n - 1)) / n
    for i, (num, lab, col) in enumerate(items):
        x = l + i * (cw + gap)
        card = rect(s, x, t, cw, 1.35, RGBColor(0xF6, 0xF8, 0xFB))
        card.line.color.rgb = RGBColor(0xD8, 0xE0, 0xEA); card.line.width = Pt(0.75)
        tf = tb(s, x + 0.12, t + 0.14, cw - 0.24, 0.7, MSO_ANCHOR.MIDDLE)
        para(tf, num, 30, col, bold=True, align=PP_ALIGN.CENTER, first=True, space=0)
        tf2 = tb(s, x + 0.12, t + 0.82, cw - 0.24, 0.48, MSO_ANCHOR.TOP)
        para(tf2, lab, 11, GRAY, align=PP_ALIGN.CENTER, first=True, space=0)


# ══════════════════════════ SLIDE 1 — TITLE ══════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
s.shapes.add_picture(str(TPL / "dots.png"), I(8.33), 0, height=I(SH))
s.shapes.add_picture(str(TPL / "logo.png"), I(0.4), I(0.35), width=I(3.6))
tf = tb(s, 0.55, 2.5, 7.7, 2.0)
para(tf, "HỆ THỐNG LIÊN CAMERA AI", 30, RED, bold=True, first=True, space=2)
para(tf, "THEO DÕI ĐỐI TƯỢNG VÀ DỰ ĐOÁN HƯỚNG ĐI", 30, RED, bold=True, space=0)
tf2 = tb(s, 0.57, 4.35, 8.4, 0.5)
para(tf2, "Phát hiện – Theo dõi – Dự đoán hướng đi trên video CCTV giao thông thật", 15, NAVY, italic=True, first=True, space=0)
tf3 = tb(s, 0.57, 5.25, 8.4, 1.6)
para(tf3, "Sinh viên: Phan Sỹ Hùng   ·   MSSV: 20225631", 14, DARK, bold=True, first=True, space=3)
para(tf3, "Chương trình: Công nghệ thông tin Việt-Nhật", 13, DARK, space=3)
para(tf3, "Giảng viên hướng dẫn: ThS. Vũ Đức Vượng", 13, DARK, space=3)
para(tf3, "Khoa Học Máy Tính · Trường CNTT & Truyền thông · Hà Nội, 07/2026", 12.5, GRAY, space=0)

# ══════════════════════════ SLIDE 2 — OUTLINE ══════════════════════════
s = content("Nội dung trình bày")
bullets(s, [
    "1.  Đặt vấn đề & động lực",
    "2.  Mục tiêu, phạm vi và hướng giải quyết",
    "3.  Kiến trúc hệ thống & phương pháp",
    "4.  Dữ liệu và quy trình huấn luyện",
    "5.  Kết quả thực nghiệm (trên dữ liệu thật)",
    "6.  Thảo luận – đánh giá trung thực",
    "7.  Kết luận, đóng góp & hướng phát triển",
], 1.1, 2.7, 10.5, 4.2, size=20, space=14)

# ══════════════════════════ SLIDE 3 — ĐẶT VẤN ĐỀ ══════════════════════════
s = content("1. Đặt vấn đề", "Giám sát giao thông bằng camera tại đô thị Việt Nam")
bullets(s, [
    "Camera giám sát tăng nhanh, nhưng phần lớn vận hành THỤ ĐỘNG: ghi hình + người trực quan sát.",
    "Ba hạn chế cốt lõi:",
    ("Không phát hiện vi phạm theo thời gian thực — chủ yếu xem lại sau khi đã xảy ra", 1),
    ("Mỗi camera độc lập, không có \"trí nhớ\" về phương tiện khi rời khung hình", 1),
    ("Chỉ phản ứng, không dự đoán được phương tiện sẽ đi hướng nào tiếp theo", 1),
    "Đặc thù VN: mật độ xe máy áp đảo, vật thể nhỏ, che khuất nhiều → khó cho các mô hình sẵn có.",
], 0.7, 2.75, 6.6, 4.2, size=15.5, space=9)
pic_fit(s, FIG / "02_cctv_pho_hue_sample_frame.jpg", 7.6, 2.75, 5.2, 3.5,
        caption="Khung hình CCTV thật — nút Phố Huế, Hà Nội")

# ══════════════════════════ SLIDE 4 — MỤC TIÊU ══════════════════════════
s = content("2. Mục tiêu, phạm vi & hướng giải quyết")
bullets(s, [
    "Mục tiêu: hệ thống giám sát AI đa camera, vận hành trên video CCTV THẬT của một tuyến nhiều giao lộ ở Hà Nội.",
    "Ba chức năng chính:",
    ("Phát hiện + theo dõi phương tiện; phán đoán vi phạm (vượt đèn đỏ/vạch dừng, sai làn)", 1),
    ("Đánh dấu đối tượng nghi vấn kèm lưu biển số khi đọc được", 1),
    ("Dự đoán HƯỚNG DI CHUYỂN của phương tiện, hiển thị trên bản đồ thật", 1),
    "Khoảng trống: chưa giải pháp nào (học thuật/thương mại) gộp cả 3 năng lực trên dữ liệu giao thông VN.",
    "So với công trình gần nhất (Chau et al., ICIT 2023 — dự đoán ngắn hạn, 1 camera): đề tài đi tới dự đoán hướng đi ở ngã tư.",
], 0.7, 2.75, 7.1, 4.3, size=14.5, space=7)
pic_fit(s, UML / "usecase_tong_quat.png", 8.05, 2.75, 4.8, 3.6,
        caption="Sơ đồ use case tổng quát")

# ══════════════════════════ SLIDE 5 — KIẾN TRÚC ══════════════════════════
s = content("3. Kiến trúc hệ thống & pipeline xử lý")
bullets(s, [
    "Pipeline nhiều mô hình chuyên biệt nối tiếp (thay vì 1 mô hình end-to-end) → mỗi thành phần tinh chỉnh riêng trên dữ liệu VN.",
    "Bốn thành phần chính:",
    ("① Phát hiện — YOLO11s_vn (5 lớp giao thông)", 1),
    ("② Theo dõi 1 camera — ByteTrack (Kalman + Hungarian)", 1),
    ("③ Dự đoán quỹ đạo ngắn hạn — Kalman Filter", 1),
    ("④ Dự đoán hướng đi — Goal Classifier (đóng góp chính)", 1),
    "Bổ trợ: phát hiện vi phạm theo luật + nhận dạng biển số; hiển thị dashboard & bản đồ.",
], 0.7, 2.75, 6.6, 4.3, size=14.5, space=7)
pic_fit(s, UML / "component_kien_truc_tong_the.png", 7.55, 2.7, 5.3, 3.65,
        caption="Kiến trúc tổng thể của hệ thống")

# ══════════════════════════ SLIDE 6 — YOLO ══════════════════════════
s = content("3.1. Phát hiện đối tượng — YOLO11s_vn")
bullets(s, [
    "YOLO11s fine-tune cho giao thông VN, 5 lớp: người, ô tô, xe máy, xe buýt, xe tải.",
    "Chọn bản 's' để cân bằng tốc độ/độ chính xác cho triển khai đa camera thời gian thực.",
    "Suy luận theo lô (batch) nhiều camera + NMS; đầu ra là hộp bao + lớp + độ tin cậy.",
    "Thách thức đã xử lý: xe máy nhỏ, mật độ cao, che khuất (mục 3.1 báo cáo).",
], 0.7, 2.75, 6.5, 4.0, size=15, space=9)
pic_fit(s, FIG / "05_detection_yolo11s_vn_predictions.jpg", 7.5, 2.75, 5.35, 3.6,
        caption="Kết quả phát hiện trên khung CCTV thật")

# ══════════════════════════ SLIDE 7 — ByteTrack ══════════════════════════
s = content("3.2. Theo dõi đa đối tượng trong 1 camera — ByteTrack")
bullets(s, [
    "Mô hình tracking-by-detection: gán ID nhất quán cho từng phương tiện qua các khung.",
    "Bộ lọc Kalman dự đoán vị trí + ghép cặp Hungarian; ngưỡng kép (high/low) giữ được vật độ tin cậy thấp.",
    "frame_rate đặt đúng theo FPS thật của từng camera (đã sửa lỗi hardcode → ảnh hưởng track_buffer).",
    "Nền tảng cho quỹ đạo & dự đoán hướng đi phía sau.",
], 0.7, 2.75, 6.5, 4.0, size=15, space=9)
pic_fit(s, FIG / "20_demo_realrun_detect_track_1.jpg", 7.5, 2.75, 5.35, 3.6,
        caption="Theo dõi phương tiện kèm track ID trên video thật")

# ══════════════════════════ SLIDE 8 — Goal Classifier ══════════════════════════
s = content("3.3. Dự đoán hướng đi — Goal Classifier", "Đóng góp chính của đồ án")
bullets(s, [
    "Bài toán: từ quỹ đạo theo dõi, phân loại hướng đi tại ngã tư: đi thẳng / rẽ trái / rẽ phải / quay đầu.",
    "33 đặc trưng chuyển động (tốc độ, biến thiên góc hướng, gia tốc ngang, kích thước hộp…).",
    "Mô hình: Gradient Boosting + hiệu chỉnh xác suất Platt (CalibratedClassifierCV).",
    "Đo ECE để mức tin cậy hiển thị phản ánh đúng tần suất đúng thực tế.",
    "Kết quả dự đoán được chiếu lên bản đồ đường thật (OpenStreetMap).",
], 0.7, 2.75, 6.7, 4.2, size=15, space=8)
pic_fit(s, FIG / "04_pho_hue_so_do_giao_lo.png", 7.75, 2.75, 5.1, 3.55,
        caption="Đồ thị giao lộ thật (OSM): 101 nút / 742 đoạn")

# ══════════════════════════ SLIDE 9 — Vi phạm + biển số ══════════════════════════
s = content("3.4. Phát hiện vi phạm & nhận dạng biển số")
bullets(s, [
    "Suy diễn theo LUẬT trên vị trí/quỹ đạo phương tiện trong khung cảnh:",
    ("Vượt đèn đỏ / vượt vạch dừng khi đèn đỏ", 1),
    ("Đi sai làn (vùng làn cấu hình theo từng camera)", 1),
    "Vùng quan tâm (ROI: vạch dừng, làn) cấu hình trực tiếp trên giao diện, đẩy nóng không cần khởi động lại.",
    "Nhận dạng biển số 2 giai đoạn (phát hiện → OCR) hỗ trợ lưu/truy xuất thông tin phương tiện nghi vấn.",
], 0.7, 2.75, 6.7, 4.0, size=15, space=9)
pic_fit(s, FIG / "22_demo_realrun_lane_violation_1.jpg", 7.75, 2.75, 5.1, 3.6,
        caption="Phát hiện đi sai làn (xe vi phạm tô đỏ)")

# ══════════════════════════ SLIDE 10 — DỮ LIỆU ══════════════════════════
s = content("4. Dữ liệu & quy trình huấn luyện", "Toàn bộ trên dữ liệu CCTV giao thông thật Việt Nam")
bullets(s, [
    "Rào cản: không có bộ dữ liệu phát hiện đủ lớn, đúng đặc thù giao thông VN.",
    "Giải pháp: TỰ ĐỘNG GÁN NHÃN (auto-label) bằng mô hình \"giáo viên\" mạnh + kiểm định mẫu thủ công (>85% hợp lệ).",
    "Tập phát hiện: 15.957 ảnh / 623.329 khung bao (hợp nhất ~124 clip CCTV).",
    "Dữ liệu Goal Classifier: 6.959 track → 5.332 huấn luyện / 1.334 kiểm thử (tách theo track).",
    "Bản đồ: OpenStreetMap khu Phố Huế – Trần Khát Chân (101 nút / 742 đoạn).",
], 0.7, 2.75, 7.0, 4.2, size=15, space=9)
pic_fit(s, FIG / "01_ban_do_hanoi_district3_topdown.png", 8.0, 2.75, 4.85, 3.6,
        caption="Khu vực thu thập dữ liệu (Hà Nội)")

# ══════════════════════════ SLIDE 11 — KẾT QUẢ YOLO ══════════════════════════
s = content("5.1. Kết quả — Phát hiện (YOLO11s_vn)")
kpi(s, [("89,5%", "mAP@0.5 (báo cáo)", NAVY),
        ("93,25%", "mAP@0.5 — nhãn NGƯỜI kiểm", GREEN),
        ("94,6%", "Precision", NAVY),
        ("83,3%", "Recall", RED)], 0.7, 2.5, 12.0)
pic_fit(s, AS / "yolo11s_vn_manual_xadan" / "confusion_matrix_normalized.png", 0.7, 3.95, 5.9, 2.4,
        caption="Ma trận nhầm lẫn (đối chiếu nhãn người)")
pic_fit(s, AS / "yolo11s_vn_xadan_qual" / "frame_002_src0026.jpg", 6.9, 3.95, 5.9, 2.4,
        caption="Phát hiện trên nút Xã Đàn (nhãn người kiểm)")

# ══════════════════════════ SLIDE 12 — KẾT QUẢ GOAL ══════════════════════════
s = content("5.2. Kết quả — Dự đoán hướng đi (Goal Classifier)")
kpi(s, [("60,12%", "Accuracy (vs nhãn auto)", NAVY),
        ("65,0%", "Accuracy — nhãn NGƯỜI (CI 58–71%)", GREEN),
        ("54,18%", "Balanced accuracy", NAVY),
        ("0,101", "ECE (hiệu chỉnh)", RED)], 0.7, 2.5, 12.0)
pic_fit(s, AS / "goal_classifier_verified" / "confusion_matrix_human.png", 0.7, 3.95, 5.9, 2.4,
        caption="Ma trận nhầm lẫn — model vs nhãn NGƯỜI (n=200)")
pic_fit(s, AS / "goal_classifier_real" / "reliability_diagram.png", 6.9, 3.95, 5.9, 2.4,
        caption="Reliability diagram (ECE = 0,101)")

# ══════════════════════════ SLIDE 13 — DEMO 1 ══════════════════════════
s = content("5.3. Sản phẩm chạy thật — Phát hiện, theo dõi & vi phạm")
pic_fit(s, FIG / "21_demo_realrun_detect_track_2.jpg", 0.7, 2.65, 5.9, 3.55,
        caption="UC1 — Phát hiện & theo dõi (track ID) trên CCTV Phố Huế")
pic_fit(s, FIG / "23_demo_realrun_lane_violation_2.jpg", 6.9, 2.65, 5.9, 3.55,
        caption="UC2 — Phát hiện vi phạm đi sai làn")

# ══════════════════════════ SLIDE 14 — DEMO 2 ══════════════════════════
s = content("5.4. Sản phẩm chạy thật — Dashboard & bản đồ")
pic_fit(s, SHOT / "14_ui_dashboard_camera_grid.png", 0.7, 2.65, 5.9, 3.55,
        caption="Dashboard giám sát đa camera thời gian thực")
pic_fit(s, SHOT / "15_ui_mapview_route_prediction.png", 6.9, 2.65, 5.9, 3.55,
        caption="UC3 — Hướng đi dự đoán hiển thị trên bản đồ thật")

# ══════════════════════════ SLIDE 15 — TRACKING/PP ĐÁNH GIÁ ══════════════════════════
s = content("5.5. Đánh giá theo dõi & phương pháp đo trung thực")
bullets(s, [
    "MOTA/IDF1 cần ground-truth track đồng bộ từng khung — video thật không có sẵn.",
    "GT bán tự động (tái dùng vết ByteTrack + sửa ID) → FP≈FN≈0 theo cấu trúc.",
    "→ MOTA (99,5%) / IDF1 (100,8%) là ARTIFACT; chỉ số lần ĐỔI ID mới có ý nghĩa.",
    "Bài học đo lường: xác thực với nhãn NGƯỜI, báo khoảng tin cậy (Wilson/bootstrap) + ECE, thừa nhận giới hạn thay vì che giấu.",
], 0.7, 2.75, 8.2, 4.0, size=15.5, space=11)
pic_fit(s, FIG / "10_tracking_mota_idf1.png", 9.0, 2.9, 3.85, 3.4,
        caption="Chỉ số MOTA/IDF1")

# ══════════════════════════ SLIDE 16 — THẢO LUẬN ══════════════════════════
s = content("6. Thảo luận — điểm mạnh & hạn chế")
bullets(s, [
    "Điểm mạnh:",
    ("Vận hành thật trên CCTV VN; xác thực bằng nhãn NGƯỜI, có khoảng tin cậy", 1),
    ("Nhãn người (65,0%) > số cũ (60,12%) → nhãn auto nhiễu 24,5% đã DÌM điểm; số báo cáo là thận trọng", 1),
    ("Phát hiện & sửa lỗi ẩn (lệch horizon train/inference; frame_rate; đọc ảnh path Unicode)", 1),
    "Hạn chế (thừa nhận):",
    ("mAP 93,25% là IN-DOMAIN (Xã Đàn có trong train) — chưa đo tổng quát hóa OOD", 1),
    ("Lớp \"quay đầu\" rất ít mẫu; recall xe máy ~79% (bỏ sót vật nhỏ)", 1),
    ("Tracking chỉ dùng IDF1/số đổi ID; MOTA chưa đo độc lập được", 1),
], 0.7, 2.7, 12.0, 4.3, size=14.5, space=6)

# ══════════════════════════ SLIDE 17 — KẾT LUẬN ══════════════════════════
s = content("7. Kết luận")
bullets(s, [
    "Xây dựng hoàn chỉnh pipeline giám sát giao thông AI 4 thành phần, vận hành thật trên CCTV Hà Nội.",
    "Phát hiện đạt mAP@0.5 = 89,5% (93,25% trên nhãn người kiểm, in-domain).",
    "Goal Classifier — đóng góp chính — dự đoán hướng đi đạt 65,0% (nhãn người, CI 58–71%), có hiệu chỉnh xác suất.",
    "Tích hợp phát hiện vi phạm, nhận dạng biển số, dashboard & bản đồ dự đoán hướng đi.",
    "Xuyên suốt: phương pháp đánh giá TRUNG THỰC — xác thực bằng nhãn người, báo khoảng tin cậy, thừa nhận giới hạn.",
], 0.75, 2.75, 12.0, 4.2, size=16, space=12)

# ══════════════════════════ SLIDE 18 — ĐÓNG GÓP ══════════════════════════
s = content("Đóng góp nổi bật")
bullets(s, [
    "Quy trình TỰ ĐỘNG GÁN NHÃN + kiểm định mẫu → vượt rào cản thiếu dữ liệu phát hiện giao thông VN.",
    "Goal Classifier dự đoán hướng đi trên dữ liệu thật: so sánh kiến trúc, hiệu chỉnh xác suất (ECE), khớp horizon train/inference.",
    "Bộ công cụ ĐÁNH GIÁ trung thực: hand-verify (nhãn người + khoảng tin cậy), quy trình đo MOTA/IDF1 và chỉ rõ giới hạn của nó.",
    "Hệ thống tích hợp đầy đủ: AI pipeline + backend + dashboard + bản đồ, đóng gói chạy demo bằng 1 lệnh.",
], 0.75, 2.75, 12.0, 4.2, size=16, space=12)

# ══════════════════════════ SLIDE 19 — HƯỚNG PHÁT TRIỂN ══════════════════════════
s = content("Hạn chế còn tồn tại & hướng phát triển")
bullets(s, [
    "Xây tập kiểm thử phát hiện ĐỘC LẬP (khác miền/OOD) để đo tổng quát hóa thật.",
    "Bổ sung dữ liệu & phương pháp cho lớp \"quay đầu\" và Re-ID xe máy (đặc thù VN, chưa có dataset công khai).",
    "Hoàn tất ground-truth theo dõi độc lập để báo MOTA/IDF1 đúng bản chất.",
    "Wire hoàn chỉnh nhận dạng biển số vào luồng vận hành; mở rộng theo dõi xuyên nhiều camera.",
    "Đo hiệu năng thực tế trên cấu hình 6 camera đồng thời.",
], 0.75, 2.75, 12.0, 4.2, size=16, space=11)

# ══════════════════════════ SLIDE 20 — THANK YOU ══════════════════════════
s = slide(TPL / "pg-13.png")
tf = tb(s, 7.0, 5.15, 5.6, 1.2)
para(tf, "Phan Sỹ Hùng — ThS. Vũ Đức Vượng (GVHD)", 15, NAVY, bold=True, first=True, space=4)
para(tf, "Trường CNTT & Truyền thông · ĐHBK Hà Nội · 07/2026", 12.5, GRAY, space=0)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("Đã tạo:", OUT, "| số slide:", len(prs.slides._sldIdLst))
