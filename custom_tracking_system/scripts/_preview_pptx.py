# -*- coding: utf-8 -*-
"""Render proxy các slide từ chính file .pptx (vị trí+text+ảnh thật) để soi layout."""
import io, sys
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parents[2]
PPTX = ROOT / "Bao_cao_do_an_tot_nghiep" / "slide_bao_cao.pptx"
OUT = ROOT / "component_for_final" / "hust_tpl" / "_prev"; OUT.mkdir(exist_ok=True, parents=True)
DPI = 110
ARIAL = "C:/Windows/Fonts/arial.ttf"; ARIALBD = "C:/Windows/Fonts/arialbd.ttf"
want = [int(x) for x in sys.argv[1:]] or [1, 3, 5, 11, 12, 16, 20]


def px(emu): return int(Emu(emu).inches * DPI)
def font(sz, bold): return ImageFont.truetype(ARIALBD if bold else ARIAL, max(8, int(sz * DPI / 72)))


def wrap(draw, text, fnt, maxw):
    words = text.split(); lines = []; cur = ""
    for w in words:
        t = (cur + " " + w).strip()
        if draw.textlength(t, font=fnt) <= maxw or not cur:
            cur = t
        else:
            lines.append(cur); cur = w
    if cur: lines.append(cur)
    return lines


prs = Presentation(str(PPTX))
W = px(prs.slide_width); H = px(prs.slide_height)
slides = list(prs.slides)
for idx in want:
    s = slides[idx - 1]
    canvas = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(canvas)
    for sh in s.shapes:
        if sh.shape_type == 13 or sh.name.startswith("Picture"):  # picture
            try:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                im = im.resize((max(1, px(sh.width)), max(1, px(sh.height))))
                canvas.paste(im, (px(sh.left), px(sh.top)))
            except Exception as e:
                d.rectangle([px(sh.left), px(sh.top), px(sh.left)+px(sh.width), px(sh.top)+px(sh.height)], outline="red")
        elif sh.has_text_frame:
            # fill (for rect shapes)
            try:
                if sh.fill.type is not None and sh.fill.fore_color and sh.fill.type == 1:
                    c = sh.fill.fore_color.rgb
                    d.rectangle([px(sh.left), px(sh.top), px(sh.left)+px(sh.width), px(sh.top)+px(sh.height)],
                                fill=(c[0], c[1], c[2]))
            except Exception:
                pass
            y = px(sh.top)
            for p in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in p.runs)
                if not txt:
                    y += int(10 * DPI / 72); continue
                r0 = p.runs[0]
                sz = r0.font.size.pt if r0.font.size else 14
                bold = bool(r0.font.bold)
                col = (0, 0, 0)
                if r0.font.color and r0.font.color.type is not None:
                    try: col = tuple(r0.font.color.rgb)
                    except Exception: pass
                fnt = font(sz, bold)
                align = p.alignment
                for line in wrap(d, txt, fnt, px(sh.width)):
                    w = d.textlength(line, font=fnt)
                    x = px(sh.left)
                    if str(align) == "CENTER (2)" or (align is not None and int(align) == 2):
                        x = px(sh.left) + (px(sh.width) - w) // 2
                    elif align is not None and int(align) == 3:
                        x = px(sh.left) + px(sh.width) - int(w)
                    d.text((x, y), line, font=fnt, fill=col)
                    y += int(sz * DPI / 72 * 1.2)
                y += int((p.space_after.pt if p.space_after else 3) * DPI / 72)
            # cảnh báo tràn
            if y > px(sh.top) + px(sh.height) + 6:
                d.rectangle([px(sh.left), px(sh.top), px(sh.left)+px(sh.width), min(H-1, y)], outline="orange", width=3)
    p = OUT / f"slide_{idx:02d}.png"
    canvas.save(p)
    print("saved", p)
